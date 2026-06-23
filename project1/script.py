from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from datetime import date
import os

# -----------------------------
# Basic setup
# -----------------------------
OUTPUT_FILE = "loan_presentation.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Theme colors
NAVY = RGBColor(18, 38, 75)
BLUE = RGBColor(46, 117, 182)
TEAL = RGBColor(30, 136, 145)
GOLD = RGBColor(204, 153, 0)
DARK = RGBColor(40, 40, 40)
GRAY = RGBColor(110, 110, 110)
LIGHT_BG = RGBColor(247, 249, 252)
WHITE = RGBColor(255, 255, 255)

# Source note used on slides
SOURCE_NOTE = "Source: streamlit_app.py and loan_amount-data-visualization.ipynb"

# -----------------------------
# Helper functions
# -----------------------------
def set_slide_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_top_bar(slide, title, subtitle=None, color=NAVY):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.82)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(11.8), Inches(0.55))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(12.0), Inches(0.35))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRAY

def add_footer(slide, page_num):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(6.86), Inches(12.4), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(225, 229, 235)
    line.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(10.5), Inches(0.25))
    p = left.text_frame.paragraphs[0]
    p.text = SOURCE_NOTE
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY

    right = slide.shapes.add_textbox(Inches(12.0), Inches(6.9), Inches(0.55), Inches(0.25))
    p2 = right.text_frame.paragraphs[0]
    p2.text = str(page_num)
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY

def add_bullets(slide, bullets, left=0.8, top=1.35, width=11.7, height=5.2, font_size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)
        p.space_before = Pt(0)
    return box

def add_two_column_bullets(slide, left_title, left_items, right_title, right_items):
    # Left column
    left_hdr = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(5.7), Inches(0.35))
    p = left_hdr.text_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE

    add_bullets(slide, left_items, left=0.78, top=1.7, width=5.8, height=4.8, font_size=18)

    # Right column
    right_hdr = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.7), Inches(0.35))
    p = right_hdr.text_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL

    add_bullets(slide, right_items, left=6.82, top=1.7, width=5.8, height=4.8, font_size=18)

def add_section_label(slide, text, left, top, width=4.0, color=GOLD):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.45)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(left), Inches(top + 0.06), Inches(width), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_table_like_box(slide, title, items, left=0.9, top=1.55, width=11.5, height=4.6):
    hdr = slide.shapes.add_textbox(Inches(left), Inches(top - 0.35), Inches(width), Inches(0.3))
    p = hdr.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = RGBColor(220, 225, 232)

    tx = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.5), Inches(height - 0.35))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)

# -----------------------------
# Slide 1: Title
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)

# Accent blocks
for x, color in [(0.55, BLUE), (0.95, TEAL), (1.35, GOLD)]:
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(0.55), Inches(0.18), Inches(5.95))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

title = slide.shapes.add_textbox(Inches(1.8), Inches(1.2), Inches(10.8), Inches(1.4))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Smart Loan Application Portal"
p.font.size = Pt(34)
p.font.bold = True
p.font.color.rgb = WHITE

sub = slide.shapes.add_textbox(Inches(1.82), Inches(2.15), Inches(10.6), Inches(0.8))
p2 = sub.text_frame.paragraphs[0]
p2.text = "Machine learning loan amount recommendation system"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(220, 230, 240)

details = slide.shapes.add_textbox(Inches(1.82), Inches(3.0), Inches(10.8), Inches(1.2))
tf3 = details.text_frame
for i, txt in enumerate([
    "Built from the attached Streamlit application and data visualization notebook.",
    "Covers data preparation, model training, and a user-facing prediction flow."
]):
    p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
    p.text = txt
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(225, 235, 245)
    p.space_after = Pt(8)

tag = slide.shapes.add_textbox(Inches(1.82), Inches(5.55), Inches(6.5), Inches(0.4))
p4 = tag.text_frame.paragraphs[0]
p4.text = "Prepared for MALHAR JOJARE"
p4.font.size = Pt(14)
p4.font.bold = True
p4.font.color.rgb = RGBColor(235, 240, 245)

add_footer(slide, 1)

# -----------------------------
# Slide 2: Project Goal
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide, "Project Goal", "What the solution is designed to do")
add_section_label(slide, "BUSINESS PROBLEM", 0.85, 1.35, 2.7, BLUE)
add_bullets(slide, [
    "Predict a recommended loan amount from borrower financial and credit profile data.",
    "Support lending decisions with a data-driven regression model.",
    "Turn the model into a simple application that non-technical users can run."
], left=0.85, top=1.95, width=11.8, height=3.2, font_size=22)
add_footer(slide, 2)

# -----------------------------
# Slide 3: Dataset Overview
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide, "Dataset Overview", "Summary from the notebook")
add_two_column_bullets(
    slide,
    "Dataset Facts",
    [
        "15,000 records.",
        "17 columns.",
        "Finance domain.",
        "Regression problem."
    ],
    "Target and Inputs",
    [
        "Target variable: recommendedloanamount.",
        "Inputs include age, income, credit score, debt ratio, savings, and loan term.",
        "Also includes property ownership, defaults, and employment details."
    ]
)
add_footer(slide, 3)

# -----------------------------
# Slide 4: Data Preparation
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide, "Data Preparation", "Cleaning and validation steps")
add_two_column_bullets(
    slide,
    "Notebook Cleaning",
    [
        "Explores missing values and outliers.",
        "Handles unusual savings values using income-based replacement logic.",
        "Uses grouped analysis by years employed."
    ],
    "App Validation",
    [
        "Validates age is at least 18.",
        "Checks email format before OTP flow.",
        "Validates credit score range and required fields."
    ]
)
add_footer(slide, 4)

# -----------------------------
# Slide 5: Exploratory Analysis
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide, "Exploratory Analysis", "How the notebook examined the data")
add_bullets(slide, [
    "The notebook uses SQL-style queries and data visualization to inspect borrower patterns.",
    "It reviews age, years employed, annual income, savings balance, debt, and credit features.",
    "The analysis helps reveal which borrower attributes are likely to influence the loan amount."
], left=0.85, top=1.6, width=11.8, height=3.5, font_size=22)
add_footer(slide, 5)

# -----------------------------
# Slide 6: Modeling Approach
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide, "Modeling Approach", "Training and model selection")
add_two_column_bullets(
    slide,
    "Models Tested",
    [
        "Linear Regression.",
        "Random Forest Regressor.",
        "XGBoost Regressor."
    ],
    "Final Choice",
    [
        "XGBoost is selected as the final model.",
        "The trained model is saved as loanamountxgb.pkl.",
        "The Streamlit app loads the persisted model for predictions."
    ]
)
add_footer(slide, 6)

# -----------------------------
# Slide 7: Performance
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide, "Performance", "Reported experimental result")
add_table_like_box(
    slide,
    "Notebook Result",
    [
        "The notebook reports very strong performance for the final XGBoost model.",
        "Test R2 is about 0.9985.",
        "This is a notebook result and should be presented as the experimental outcome of the training run."
    ],
    left=1.0, top=1.7, width=11.3, height=3.5
)
add_footer(slide, 7)

# -----------------------------
# Slide 8: App Workflow
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide, "Streamlit App Workflow", "How the user experience is structured")
add_bullets(slide, [
    "Step 1: Personal information.",
    "Step 2: Email verification with OTP.",
    "Step 3: Employment details.",
    "Step 4: Credit profile.",
    "Step 5: Financials.",
    "Step 6: Loan details.",
    "Step 7: Review and prediction."
], left=0.9, top=1.45, width=6.2, height=4.8, font_size=18)

add_section_label(slide, "SESSION STATE + MODEL PREDICTION", 7.0, 1.55, 4.5, TEAL)
add_bullets(slide, [
    "The app stores each step in session state.",
    "Inputs are aligned to the model’s feature order.",
    "The prediction is shown as a recommended loan amount."
], left=7.0, top=2.15, width=5.6, height=3.3, font_size=18)
add_footer(slide, 8)

# -----------------------------
# Slide 9: Key Takeaways
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide, "Key Takeaways", "What this project demonstrates")
add_two_column_bullets(
    slide,
    "Analytics Side",
    [
        "Data cleaning and outlier handling.",
        "Feature understanding through EDA.",
        "Model comparison and persistence."
    ],
    "Product Side",
    [
        "Multi-step user experience.",
        "Input validation and session handling.",
        "End-to-end loan recommendation portal."
    ]
)
add_footer(slide, 9)

# -----------------------------
# Slide 10: Closing
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
box = slide.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.3), Inches(3.2))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "From notebook to product"
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = WHITE

for i, txt in enumerate([
    "The notebook builds the modeling foundation.",
    "The Streamlit app turns it into a usable loan recommendation experience.",
    "Together, they show a complete ML workflow from analysis to deployment."
]):
    p = tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(225, 235, 245)
    p.space_after = Pt(10)

end = slide.shapes.add_textbox(Inches(1.0), Inches(5.45), Inches(11.3), Inches(0.5))
p2 = end.text_frame.paragraphs[0]
p2.text = "Thank you"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = RGBColor(235, 240, 245)

add_footer(slide, 10)

# -----------------------------
# Save file
# -----------------------------
prs.save(OUTPUT_FILE)
print(f"Created {OUTPUT_FILE}")